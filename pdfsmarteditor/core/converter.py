import os
import shutil
from typing import List

from pdf2docx import Converter


def _require_dependency(command: str, friendly_name: str):
    """Ensure an external binary exists before running a conversion."""
    if not shutil.which(command):
        raise RuntimeError(
            f"{friendly_name} is required for this operation but '{command}' was not found on PATH."
        )


class PDFConverter:
    def __init__(self):
        pass

    def pdf_to_word(self, pdf_path: str, output_path: str):
        """
        Convert PDF to Word (DOCX).
        """
        cv = Converter(pdf_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()

    def pdf_to_ppt(self, pdf_path: str, output_path: str):
        """
        Convert PDF to PowerPoint (PPTX) by extracting text and images.
        """
        import io

        import fitz
        from pptx import Presentation
        from pptx.util import Pt

        doc = fitz.open(pdf_path)
        prs = Presentation()

        # Remove default slides if any
        for i in range(len(prs.slides) - 1, -1, -1):
            rId = prs.slides._sle[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sle[i]

        for page in doc:
            # Match slide size to page size (approximate)
            page_width = page.rect.width
            page_height = page.rect.height
            prs.slide_width = Pt(page_width)
            prs.slide_height = Pt(page_height)

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # Extract text blocks
            blocks = page.get_text("dict")["blocks"]

            for b in blocks:
                if b["type"] == 0:  # Text
                    # Bounding box for the whole block
                    bbox = b["bbox"]
                    x, y = bbox[0], bbox[1]
                    w = bbox[2] - bbox[0]
                    h = bbox[3] - bbox[1]

                    txBox = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
                    tf = txBox.text_frame
                    tf.word_wrap = True

                    for line in b["lines"]:
                        p = tf.add_paragraph()
                        for span in line["spans"]:
                            run = p.add_run()
                            run.text = span["text"]
                            run.font.size = Pt(span["size"])
                            # Basic font mapping attempt
                            if "bold" in span["font"].lower():
                                run.font.bold = True
                            if "italic" in span["font"].lower():
                                run.font.italic = True

            # Extract images
            image_list = page.get_images(full=True)
            for img in image_list:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                rects = page.get_image_rects(xref)
                for rect in rects:
                    image_stream = io.BytesIO(image_bytes)
                    try:
                        slide.shapes.add_picture(
                            image_stream,
                            Pt(rect.x0),
                            Pt(rect.y0),
                            width=Pt(rect.width),
                            height=Pt(rect.height),
                        )
                    except Exception:
                        continue

        prs.save(output_path)
        doc.close()

    def pdf_to_excel(self, pdf_path: str, output_path: str):
        """
        Convert PDF to Excel (XLSX) by extracting tables.
        """
        import pandas as pd
        import pdfplumber

        with pdfplumber.open(pdf_path) as pdf:
            all_tables = []
            for i, page in enumerate(pdf.pages):
                # Try multiple strategies to find tables
                found_on_page = []

                # 1. Default strategy
                ts = page.extract_tables()
                if ts:
                    found_on_page.extend(ts)

                # 2. Lines strategy (good for clear grid lines)
                ts_lines = page.extract_tables(
                    {"vertical_strategy": "lines", "horizontal_strategy": "lines"}
                )
                if ts_lines:
                    found_on_page.extend(ts_lines)

                # 3. Text strategy (good for white-space separated tables)
                if not found_on_page:
                    ts_text = page.extract_tables(
                        {"vertical_strategy": "text", "horizontal_strategy": "text"}
                    )
                    if ts_text:
                        found_on_page.extend(ts_text)

                if found_on_page:
                    for table in found_on_page:
                        if not table:
                            continue
                        # Clean table data: remove None, strip whitespace
                        cleaned_table = [
                            [
                                str(cell).strip() if cell is not None else ""
                                for cell in row
                            ]
                            for row in table
                        ]
                        # Remove completely empty rows/columns
                        df = pd.DataFrame(cleaned_table)
                        df.dropna(how="all", axis=0, inplace=True)
                        df.dropna(how="all", axis=1, inplace=True)
                        if not df.empty:
                            all_tables.append(df)
                else:
                    # Fallback: Extract text
                    text = page.extract_text()
                    if text:
                        lines = [
                            line.split() for line in text.split("\n") if line.strip()
                        ]
                        if lines:
                            all_tables.append(pd.DataFrame(lines))

            if not all_tables:
                all_tables.append(pd.DataFrame([["No tables or text found in PDF"]]))

            # Save to Excel
            with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
                for idx, df in enumerate(all_tables):
                    sheet_name = f"Table_{idx+1}"[:31]
                    df.to_excel(
                        writer, sheet_name=sheet_name, index=False, header=False
                    )

    def word_to_pdf(self, word_path: str, output_dir: str) -> str:
        """
        Convert Word (DOC/DOCX) to PDF using LibreOffice.
        Returns the path to the created PDF.
        """
        import subprocess

        _require_dependency("libreoffice", "LibreOffice")

        # LibreOffice converts to the same directory
        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            word_path,
        ]

        try:
            subprocess.run(
                cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )

            # Construct expected output path
            filename = os.path.basename(word_path)
            name_without_ext = os.path.splitext(filename)[0]
            pdf_filename = f"{name_without_ext}.pdf"
            return os.path.join(output_dir, pdf_filename)

        except subprocess.CalledProcessError as e:
            raise Exception(f"LibreOffice conversion failed: {e.stderr.decode()}")

    def ppt_to_pdf(self, ppt_path: str, output_dir: str) -> str:
        """
        Convert PowerPoint (PPT/PPTX) to PDF using LibreOffice.
        Returns the path to the created PDF.
        """
        import subprocess

        _require_dependency("libreoffice", "LibreOffice")

        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            ppt_path,
        ]

        try:
            subprocess.run(
                cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )

            filename = os.path.basename(ppt_path)
            name_without_ext = os.path.splitext(filename)[0]
            pdf_filename = f"{name_without_ext}.pdf"
            return os.path.join(output_dir, pdf_filename)

        except subprocess.CalledProcessError as e:
            raise Exception(f"LibreOffice conversion failed: {e.stderr.decode()}")

    def excel_to_pdf(self, excel_path: str, output_dir: str) -> str:
        """
        Convert Excel (XLS/XLSX) to PDF using LibreOffice.
        Returns the path to the created PDF.
        """
        import subprocess

        _require_dependency("libreoffice", "LibreOffice")

        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            excel_path,
        ]

        try:
            subprocess.run(
                cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )

            filename = os.path.basename(excel_path)
            name_without_ext = os.path.splitext(filename)[0]
            pdf_filename = f"{name_without_ext}.pdf"
            return os.path.join(output_dir, pdf_filename)

        except subprocess.CalledProcessError as e:
            raise Exception(f"LibreOffice conversion failed: {e.stderr.decode()}")

    def pdf_to_jpg(self, pdf_path: str, output_dir: str) -> List[str]:
        """
        Convert PDF pages to JPG images.
        Returns list of paths to created images.
        """
        import fitz

        doc = fitz.open(pdf_path)
        output_files = []

        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom
            output_filename = f"page_{i+1}_{os.path.basename(pdf_path)}.jpg"
            output_path = os.path.join(output_dir, output_filename)
            pix.save(output_path)
            output_files.append(output_path)

        doc.close()
        return output_files

    def jpg_to_pdf(self, image_paths: List[str], output_path: str):
        """
        Convert JPG/PNG images to a single PDF.
        """
        import fitz

        doc = fitz.open()

        for img_path in image_paths:
            img = fitz.open(img_path)
            pdfbytes = img.convert_to_pdf()
            img.close()

            imgPDF = fitz.open("pdf", pdfbytes)
            page_src = imgPDF[0]
            rect = page_src.rect

            page = doc.new_page(width=rect.width, height=rect.height)
            page.show_pdf_page(rect, imgPDF, 0)

        doc.save(output_path)
        doc.close()

    def html_to_pdf(self, html_path: str, output_path: str):
        """
        Convert HTML to PDF.
        """
        import fitz

        doc = fitz.open(html_path)
        pdfbytes = doc.convert_to_pdf()
        with open(output_path, "wb") as f:
            f.write(pdfbytes)
        doc.close()

    def pdf_to_pdfa(self, pdf_path: str, output_path: str):
        """
        Convert PDF to PDF/A using Ghostscript.
        """
        import subprocess

        _require_dependency("gs", "Ghostscript")

        # Try PDF/A-2b
        cmd = [
            "gs",
            "-dPDFA=2",
            "-dBATCH",
            "-dNOPAUSE",
            "-sColorConversionStrategy=RGB",
            "-sProcessColorModel=DeviceRGB",
            "-sDEVICE=pdfwrite",
            "-sPDFACompatibilityPolicy=1",
            f"-sOutputFile={output_path}",
            pdf_path,
        ]

        try:
            subprocess.run(
                cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )
        except subprocess.CalledProcessError as e:
            print(f"Ghostscript failed: {e.stderr.decode()}")
            # Fallback to simple pdfwrite if PDF/A fails (e.g. missing color profile)
            print("Falling back to standard PDF conversion...")
            cmd_fallback = [
                "gs",
                "-dBATCH",
                "-dNOPAUSE",
                "-sDEVICE=pdfwrite",
                f"-sOutputFile={output_path}",
                pdf_path,
            ]
            subprocess.run(
                cmd_fallback, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )

    def scan_to_pdf(
        self, image_paths: List[str], output_path: str, enhance: bool = True
    ):
        """
        Convert scanned images to PDF, optionally enhancing them.
        """
        import fitz
        from PIL import Image, ImageEnhance

        doc = fitz.open()

        for img_path in image_paths:
            if enhance:
                # Enhance image using Pillow
                try:
                    img = Image.open(img_path)
                    # Convert to grayscale
                    img = img.convert("L")
                    # Increase contrast
                    enhancer = ImageEnhance.Contrast(img)
                    img = enhancer.enhance(1.5)

                    # Save to temp file
                    temp_img_path = f"{img_path}_enhanced.jpg"
                    img.save(temp_img_path)

                    img_doc = fitz.open(temp_img_path)
                    pdfbytes = img_doc.convert_to_pdf()
                    img_pdf = fitz.open("pdf", pdfbytes)
                    doc.insert_pdf(img_pdf)
                    img_doc.close()
                    img_pdf.close()

                    # Cleanup temp file
                    if os.path.exists(temp_img_path):
                        os.remove(temp_img_path)

                except Exception as e:
                    print(f"Failed to enhance image {img_path}: {e}")
                    # Fallback to original
                    img_doc = fitz.open(img_path)
                    pdfbytes = img_doc.convert_to_pdf()
                    img_pdf = fitz.open("pdf", pdfbytes)
                    doc.insert_pdf(img_pdf)
                    img_doc.close()
                    img_pdf.close()
            else:
                img_doc = fitz.open(img_path)
                pdfbytes = img_doc.convert_to_pdf()
                img_pdf = fitz.open("pdf", pdfbytes)
                doc.insert_pdf(img_pdf)
                img_doc.close()
                img_pdf.close()

        doc.save(output_path)
        doc.close()

    def ocr_pdf(self, pdf_path: str, output_path: str, language: str = "eng"):
        """
        OCR PDF using pytesseract (Tesseract).
        """
        import io

        import fitz
        import pytesseract
        from PIL import Image

        _require_dependency("tesseract", "Tesseract OCR")

        doc = fitz.open(pdf_path)
        out_doc = fitz.open()

        for page in doc:
            # Get image from page
            pix = page.get_pixmap()
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))

            # OCR to PDF
            try:
                pdf_bytes = pytesseract.image_to_pdf_or_hocr(
                    img, extension="pdf", lang=language
                )
                img_pdf = fitz.open("pdf", pdf_bytes)
                out_doc.insert_pdf(img_pdf)
                img_pdf.close()
            except Exception as e:
                print(f"OCR failed for page: {e}")
                # Fallback: just insert original page (as image or original)
                # If we insert original page, it might not be searchable if it was image-only.
                # But better than failing.
                out_doc.insert_pdf(doc, from_page=page.number, to_page=page.number)

        out_doc.save(output_path)
        out_doc.close()
        doc.close()
