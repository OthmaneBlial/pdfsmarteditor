import os

import fitz
import pandas as pd
import pytest
from pptx import Presentation

from pdfsmarteditor.core.converter import PDFConverter


@pytest.fixture
def sample_pdf(tmp_path):
    pdf_path = tmp_path / "test.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Hello World")
    page.insert_text((50, 100), "Test Line 2")
    doc.save(pdf_path)
    doc.close()
    return str(pdf_path)


@pytest.fixture
def sample_table_pdf(tmp_path):
    pdf_path = tmp_path / "table.pdf"
    # Create a PDF with a simple table-like structure using text
    doc = fitz.open()
    page = doc.new_page()
    # Header
    page.insert_text((50, 50), "Col1")
    page.insert_text((150, 50), "Col2")
    # Row 1
    page.insert_text((50, 70), "Val1")
    page.insert_text((150, 70), "Val2")
    doc.save(pdf_path)
    doc.close()
    return str(pdf_path)


def test_pdf_to_ppt(sample_pdf, tmp_path):
    converter = PDFConverter()
    output_path = tmp_path / "output.pptx"

    converter.pdf_to_ppt(sample_pdf, str(output_path))

    assert os.path.exists(output_path)

    # Verify content
    prs = Presentation(output_path)
    assert len(prs.slides) == 1

    text_found = False
    for shape in prs.slides[0].shapes:
        if hasattr(shape, "text"):
            if "Hello World" in shape.text:
                text_found = True
                break

    assert text_found, "Text 'Hello World' not found in generated PPTX"


def test_pdf_to_excel(sample_table_pdf, tmp_path):
    converter = PDFConverter()
    output_path = tmp_path / "output.xlsx"

    converter.pdf_to_excel(sample_table_pdf, str(output_path))

    assert os.path.exists(output_path)

    # Verify content
    df = pd.read_excel(output_path, header=None)
    # Since we used a simple text PDF, it might fall back to text extraction or find a table depending on pdfplumber
    # We just want to ensure it produced something valid and contains our data

    content = df.to_string()
    assert "Val1" in content or "Val1" in str(df.values)
    assert "Val2" in content or "Val2" in str(df.values)


def test_pdf_to_word(sample_pdf, tmp_path):
    converter = PDFConverter()
    output_path = tmp_path / "output.docx"

    converter.pdf_to_word(sample_pdf, str(output_path))

    assert os.path.exists(output_path)
    # We can't easily verify docx content without python-docx, but existence is a good start.
    # If we wanted to be thorough:
    # from docx import Document
    # doc = Document(output_path)
    # full_text = []
    # for para in doc.paragraphs:
    #     full_text.append(para.text)
    # assert "Hello World" in '\n'.join(full_text)


def test_scan_to_pdf(tmp_path):
    converter = PDFConverter()
    output_path = tmp_path / "scanned.pdf"

    # Create dummy images
    img_paths = []
    for i in range(2):
        img_path = tmp_path / f"scan_{i}.jpg"
        pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 100, 100), 0)
        pix.save(str(img_path))
        img_paths.append(str(img_path))

    converter.scan_to_pdf(img_paths, str(output_path), enhance=False)

    assert os.path.exists(output_path)
    doc = fitz.open(output_path)
    assert len(doc) == 2
    doc.close()


def test_ocr_pdf(sample_pdf, tmp_path):
    # This test might fail if tesseract is not installed or configured.
    # We should skip if tesseract is missing or mock it.
    # For now, let's assume we can mock the dependency check or just try.

    converter = PDFConverter()
    output_path = tmp_path / "ocr.pdf"

    # We need to mock _require_dependency if we don't want to rely on system tools
    # But let's try to run it if possible, or catch the error.

    try:
        converter.ocr_pdf(sample_pdf, str(output_path))
        assert os.path.exists(output_path)
    except RuntimeError as e:
        pytest.skip(f"Skipping OCR test: {e}")
    except Exception as e:
        # If tesseract fails for other reasons (e.g. language data missing)
        pytest.skip(f"OCR failed: {e}")
